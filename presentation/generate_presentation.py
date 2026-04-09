"""
Generate 5-slide PowerPoint presentation for Procurement Tracker project.
Run: py generate_presentation.py
"""

import os
import urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== COLOR PALETTE =====
CLR_TITLE = RGBColor(0x1E, 0x3A, 0x5F)
CLR_ACCENT = RGBColor(0x3B, 0x82, 0xF6)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
CLR_TEXT = RGBColor(0x33, 0x41, 0x55)
CLR_SUBTEXT = RGBColor(0x64, 0x74, 0x8B)
CLR_GREEN = RGBColor(0x22, 0xC5, 0x5E)
CLR_GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
CLR_RED = RGBColor(0xF8, 0x71, 0x71)
CLR_RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
CLR_CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
CLR_YELLOW = RGBColor(0xFB, 0xBC, 0x24)

# ===== URLS =====
GITHUB_URL = "https://github.com/MarikSH/Hakaton_SET_ZakupBot"
DEPLOY_URL = ""  # Will be filled after deployment

# ===== QR CODE GENERATION =====
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
QR_GITHUB_PATH = os.path.join(SCRIPT_DIR, "qr_github.png")
QR_DEPLOY_PATH = os.path.join(SCRIPT_DIR, "qr_deploy.png")

def download_qr(url, save_path, size=200):
    """Download QR code image from qrserver.com API."""
    qr_url = f"https://api.qrserver.com/v1/create-qr-code/?size={size}x{size}&data={urllib.parse.quote(url)}&color=3b82f6&bgcolor=ffffff"
    try:
        urllib.request.urlretrieve(qr_url, save_path)
        print(f"QR code saved to: {save_path}")
        return True
    except Exception as e:
        print(f"Failed to download QR code: {e}")
        return False

import urllib.parse
download_qr(GITHUB_URL, QR_GITHUB_PATH)
# Placeholder QR for deploy
download_qr("https://github.com/MarikSH/Hakaton_SET_ZakupBot", QR_DEPLOY_PATH)

# ===== HELPER FUNCTIONS =====

def add_bg(slide, color=CLR_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=CLR_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, bold=False, color=CLR_TEXT, alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def add_rich_tf(tf, parts):
    """Add rich formatted text to a text frame. parts = list of (text, size, bold, color)."""
    for i, (text, size, bold, color) in enumerate(parts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return tf


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, CLR_WHITE)
add_shape(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_textbox(slide1, Inches(6.2), Inches(0.8), Inches(1), Inches(1), "📦",
            font_size=52, alignment=PP_ALIGN.CENTER)

add_textbox(slide1, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
            "Procurement Tracker", font_size=44, bold=True, color=CLR_ACCENT,
            alignment=PP_ALIGN.CENTER, font_name='Calibri Light')

add_textbox(slide1, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.8),
            "Mobile web app for marketplace sellers — tracking procurement costs with per-unit cost analysis and budget control",
            font_size=20, color=CLR_SUBTEXT, alignment=PP_ALIGN.CENTER)

add_shape(slide1, Inches(5.5), Inches(3.9), Inches(2.3), Inches(0.03), CLR_ACCENT)

author_tf = add_textbox(slide1, Inches(1.5), Inches(4.3), Inches(10.3), Inches(2.5), "", font_size=20, color=CLR_TEXT, alignment=PP_ALIGN.CENTER)
p = author_tf.paragraphs[0]
p.text = "Marat Sharapov"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = CLR_TITLE
p.alignment = PP_ALIGN.CENTER

add_para(author_tf, "m.sharapov@innopolis.university", font_size=18, color=CLR_SUBTEXT, alignment=PP_ALIGN.CENTER)
add_para(author_tf, "Group CSE-04", font_size=18, color=CLR_SUBTEXT, alignment=PP_ALIGN.CENTER)
add_para(author_tf, "", font_size=8, alignment=PP_ALIGN.CENTER)
add_para(author_tf, "github.com/MarikSH/Hakaton_SET_ZakupBot", font_size=16,
         color=CLR_ACCENT, alignment=PP_ALIGN.CENTER, bold=True)

add_shape(slide1, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_speaker_notes(slide1,
    "Good morning/afternoon. My name is Marat Sharapov from Innopolis University, group CSE-04.\n"
    "I present Procurement Tracker — a mobile web application designed for marketplace sellers "
    "who need accurate per-unit cost calculations across all their procurement expenses.\n"
    "The full source code is available at github.com/MarikSH/Hakaton_SET_ZakupBot.")


# ============================================================
# SLIDE 2: CONTEXT
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, CLR_WHITE)
add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_textbox(slide2, Inches(0.8), Inches(0.3), Inches(5), Inches(0.7),
            "Context", font_size=32, bold=True, color=CLR_TITLE)

# End Users
user_card = add_rounded_rect(slide2, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.6),
                              RGBColor(0xEF, 0xF6, 0xFF), CLR_ACCENT)
user_tf = add_textbox(slide2, Inches(1.2), Inches(1.3), Inches(11), Inches(1.4), "", font_size=20)
p = user_tf.paragraphs[0]
p.text = "End Users"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = CLR_TITLE
add_para(user_tf, "Marketplace sellers (Wildberries, Ozon), small business owners, and entrepreneurs "
         "who procure goods wholesale and need precise per-unit cost tracking including delivery, "
         "packaging, staff, and warehouse expenses.", font_size=18, color=CLR_TEXT)

# Problem
prob_card = add_rounded_rect(slide2, Inches(0.8), Inches(3.1), Inches(5.3), Inches(2.2),
                              CLR_RED_BG, CLR_RED)
prob_tf = add_textbox(slide2, Inches(1.1), Inches(3.2), Inches(4.8), Inches(2.0), "", font_size=18)
p = prob_tf.paragraphs[0]
p.text = "Problem"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = CLR_TITLE
add_para(prob_tf, "Entrepreneurs lose money from inaccurate cost calculations", font_size=17, color=CLR_TEXT)
add_para(prob_tf, "Overhead expenses (delivery, packaging, warehouse, staff) distributed manually", font_size=17, color=CLR_TEXT)
add_para(prob_tf, "No budget control leads to overspending going unnoticed", font_size=17, color=CLR_TEXT)

# Solution
sol_card = add_rounded_rect(slide2, Inches(7.2), Inches(3.1), Inches(5.3), Inches(2.2),
                              CLR_GREEN_BG, CLR_GREEN)
sol_tf = add_textbox(slide2, Inches(7.5), Inches(3.2), Inches(4.8), Inches(2.0), "", font_size=18)
p = sol_tf.paragraphs[0]
p.text = "Solution"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = CLR_TITLE
add_para(sol_tf, "Automatic per-unit cost calculation:", font_size=17, color=CLR_TEXT, bold=True)
add_para(sol_tf, "ALL expenses total divided by goods quantity = true cost per product", font_size=17, color=CLR_TEXT)
add_para(sol_tf, "Visual budget tracking with progress bars and overspend alerts", font_size=17, color=CLR_TEXT)

# One-liner
add_textbox(slide2, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5),
            "Idea: \"Enter your procurement, add expenses in natural language — get the exact cost per unit for accurate profit analysis.\"",
            font_size=19, bold=True, color=CLR_TITLE)

add_shape(slide2, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_speaker_notes(slide2,
    "Our end users are marketplace sellers on platforms like Wildberries and Ozon — people who buy goods "
    "wholesale and resell them, plus small business owners and procurement managers.\n\n"
    "The problem: sellers often miscalculate their true per-unit cost because they don't properly "
    "distribute overhead expenses like delivery, packaging, warehouse rent, and worker salaries. "
    "This leads to incorrect pricing, underestimated costs, and lost profits.\n\n"
    "Our solution: enter your procurement, add all expenses in natural language, "
    "and the system automatically calculates the exact cost per unit of goods. "
    "Key formula: total_all_expenses / quantity_of_goods = true cost per product.")


# ============================================================
# SLIDE 3: IMPLEMENTATION
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, CLR_WHITE)
add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_textbox(slide3, Inches(0.8), Inches(0.3), Inches(6), Inches(0.7),
            "Implementation", font_size=32, bold=True, color=CLR_TITLE)

# Tech Stack
stack_tf = add_textbox(slide3, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.4), "", font_size=17)
tech_items = [
    ("Backend: ", "FastAPI (Python REST API)"),
    ("Database: ", "PostgreSQL + SQLAlchemy ORM"),
    ("Frontend: ", "React + Vite + TailwindCSS"),
    ("Features: ", "PWA + i18n (RU/EN) + PDF export + Dark theme"),
    ("Deploy: ", "Docker Compose -> Railway/Vercel/VPS"),
]
for i, (label, value) in enumerate(tech_items):
    if i == 0:
        p = stack_tf.paragraphs[0]
    else:
        p = stack_tf.add_paragraph()
        p.space_before = Pt(3)
    run1 = p.add_run()
    run1.text = label
    run1.font.size = Pt(17)
    run1.font.bold = True
    run1.font.color.rgb = CLR_TITLE
    run1.font.name = 'Calibri'
    run2 = p.add_run()
    run2.text = value
    run2.font.size = Pt(17)
    run2.font.color.rgb = CLR_TEXT
    run2.font.name = 'Calibri'

# V1 column
v1_card = add_rounded_rect(slide3, Inches(0.8), Inches(3.0), Inches(5.5), Inches(2.5), CLR_CARD_BG)
v1_tf = add_textbox(slide3, Inches(1.1), Inches(3.05), Inches(5), Inches(2.4), "", font_size=17)
p = v1_tf.paragraphs[0]
p.text = "Version 1.0 - MVP"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = CLR_ACCENT

v1_items = [
    "Core API: create procurements, calculate expenses",
    "NLP parser - 6 categories, 100+ keywords (RU/EN)",
    "Formula: price_per_unit + (overhead / qty)",
    "Basic form UI for products + free-text expenses",
    "Docker Compose one-command deployment",
]
for i, item in enumerate(v1_items):
    p = v1_tf.add_paragraph()
    p.space_before = Pt(4)
    p.text = " - " + item
    p.font.size = Pt(16)
    p.font.color.rgb = CLR_TEXT

# V2 column
v2_card = add_rounded_rect(slide3, Inches(6.8), Inches(3.0), Inches(5.5), Inches(2.5), CLR_CARD_BG)
v2_tf = add_textbox(slide3, Inches(7.1), Inches(3.05), Inches(5), Inches(2.4), "", font_size=17)
p = v2_tf.paragraphs[0]
p.text = "Version 2.0 - Product"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = CLR_GREEN

v2_items = [
    "PWA - installable on mobile devices",
    "i18n - full Russian + English localization",
    "PDF export of cost reports (jsPDF + Cyrillic)",
    "Procurement statuses: active -> completed",
    "History screen + budget control with alerts",
    "Dark theme + settings modal",
]
for i, item in enumerate(v2_items):
    p = v2_tf.add_paragraph()
    p.space_before = Pt(3)
    p.text = " - " + item
    p.font.size = Pt(16)
    p.font.color.rgb = CLR_TEXT

# TA Feedback
fb_label = add_textbox(slide3, Inches(0.8), Inches(5.7), Inches(12), Inches(0.4),
            "TA Feedback Addressed:", font_size=18, bold=True, color=CLR_GREEN)

fb_card = add_rounded_rect(slide3, Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.1), CLR_GREEN_BG, CLR_GREEN)
fb_tf = add_textbox(slide3, Inches(1.1), Inches(6.15), Inches(11.2), Inches(1.0), "", font_size=16)
fb_items = [
    "Expanded NLP parser from 5 to 6+ categories with bilingual keywords (RU/EN)",
    "Added budget control with visual progress bars and overspend alerts",
    "Improved UI: added history screen, settings modal, and dark theme support",
    "Added PDF export with Cyrillic font support for professional reporting",
]
for i, item in enumerate(fb_items):
    if i == 0:
        p = fb_tf.paragraphs[0]
    else:
        p = fb_tf.add_paragraph()
        p.space_before = Pt(2)
    p.text = " [check] " + item
    p.font.size = Pt(15)
    p.font.color.rgb = CLR_TEXT

add_shape(slide3, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_speaker_notes(slide3,
    "Tech Stack: Backend on FastAPI with PostgreSQL and SQLAlchemy ORM. "
    "Frontend uses React with Vite and TailwindCSS. "
    "PWA support, full i18n for Russian and English, PDF report export with Cyrillic fonts. "
    "Docker Compose for one-command deployment.\n\n"
    "Version 1 (MVP): Core API, NLP expense parser with 6 categories and 100+ bilingual keywords, "
    "the core per-unit formula, basic form UI, and Docker Compose.\n\n"
    "Version 2 (Product): PWA capabilities, internationalization, PDF export, "
    "procurement status tracking, history screen, dark theme, budget control with alerts.\n\n"
    "TA Feedback Addressed:\n"
    "1. Expanded NLP parser from 5 to 6+ categories with bilingual keywords\n"
    "2. Added budget control with visual progress bars and overspend alerts\n"
    "3. Improved UI with history screen, settings, and dark theme\n"
    "4. Added PDF export for professional reporting")


# ============================================================
# SLIDE 4: DEMO
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, CLR_WHITE)
add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_textbox(slide4, Inches(0.8), Inches(0.3), Inches(6), Inches(0.7),
            "Demo - Version 2.0", font_size=32, bold=True, color=CLR_TITLE)

# Video placeholder
video_box = add_rounded_rect(slide4, Inches(0.8), Inches(1.2), Inches(7.5), Inches(4.5),
                              RGBColor(0x1E, 0x29, 0x3B))
add_textbox(slide4, Inches(0.8), Inches(2.5), Inches(7.5), Inches(1),
            "", font_size=60, alignment=PP_ALIGN.CENTER, color=RGBColor(0x94, 0xA3, 0xB8))
add_textbox(slide4, Inches(0.8), Inches(2.8), Inches(7.5), Inches(0.6),
            "Pre-recorded Video Demonstration", font_size=22, bold=True,
            alignment=PP_ALIGN.CENTER, color=RGBColor(0x94, 0xA3, 0xB8))
add_textbox(slide4, Inches(0.8), Inches(3.5), Inches(7.5), Inches(1.8),
            "Insert your demo video here (mp4, up to 2 minutes with voice-over)\n\n"
            "PowerPoint: Insert -> Video -> This Device -> select demo.mp4\n\n"
            "This is the MOST IMPORTANT slide - ensure the demo shows:\n"
            "  - Procurement creation\n"
            "  - NLP expense parsing\n"
            "  - Per-unit cost calculation\n"
            "  - PDF export",
            font_size=15, alignment=PP_ALIGN.CENTER, color=RGBColor(0x64, 0x74, 0x8B))

# Demo flow
demo_card = add_rounded_rect(slide4, Inches(8.8), Inches(1.2), Inches(3.8), Inches(4.5), CLR_CARD_BG, CLR_ACCENT)
demo_tf = add_textbox(slide4, Inches(9.0), Inches(1.3), Inches(3.4), Inches(4.3), "", font_size=16)
p = demo_tf.paragraphs[0]
p.text = "Demo Flow"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = CLR_TITLE

demo_steps = [
    "1. Create Procurement",
    "   Enter name + budget (200K rub)",
    "2. Add Products",
    "   e.g., 100 units x 1,000 rub",
    "3. NLP Expense Parse",
    "   'Delivery 3000, pkg 40/unit'",
    "4. View Report",
    "   Per-unit cost = 1,070 rub/unit",
    "5. Export PDF",
    "   Download full cost report",
    "6. Complete + History",
    "   Mark done, review later",
]
for i, step in enumerate(demo_steps):
    p = demo_tf.add_paragraph()
    p.space_before = Pt(3)
    p.text = step
    p.font.size = Pt(14)
    p.font.color.rgb = CLR_TEXT
    if step.strip().startswith(("1.", "2.", "3.", "4.", "5.", "6.")):
        p.font.bold = True

# Important note
note_box = add_rounded_rect(slide4, Inches(0.8), Inches(5.9), Inches(7.5), Inches(1.3),
                             RGBColor(0xFF, 0xFB, 0xEB), CLR_YELLOW)
note_tf = add_textbox(slide4, Inches(1.1), Inches(5.95), Inches(7), Inches(1.2), "", font_size=16)
p = note_tf.paragraphs[0]
p.text = "MOST IMPORTANT SLIDE"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = CLR_TITLE
add_para(note_tf, "The demo video is the critical part of this presentation. "
         "It must show the complete V2 workflow: procurement creation, NLP parsing, "
         "per-unit calculation, and PDF export - all within 2 minutes with voice-over.", font_size=15, color=CLR_TEXT)

add_shape(slide4, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_speaker_notes(slide4,
    "This is the most important part - the live demonstration of Version 2.\n\n"
    "Demo walkthrough:\n"
    "1. Create Procurement: Enter name and budget of 200,000 rubles.\n"
    "2. Add Products: 100 units at 1,000 rubles each = 100,000 rubles.\n"
    "3. NLP Expense Parsing: Type 'delivery 3000, packaging 40 per unit' - "
    "system auto-categorizes into delivery and packaging.\n"
    "4. View Report: Per-unit cost = 1,070 rubles (base price + overhead share).\n"
    "5. Export PDF: Download professional cost report.\n"
    "6. Complete and History: Mark done, appears in history.\n\n"
    "Please replace the video placeholder with your screen recording. "
    "Max 2 minutes with voice-over narration.")


# ============================================================
# SLIDE 5: LINKS
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, CLR_WHITE)
add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_textbox(slide5, Inches(0.8), Inches(0.3), Inches(6), Inches(0.7),
            "Links and QR Codes", font_size=32, bold=True, color=CLR_TITLE)

# GitHub link card
gh_card = add_rounded_rect(slide5, Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.5), CLR_CARD_BG, CLR_ACCENT)
gh_tf = add_textbox(slide5, Inches(1.2), Inches(1.3), Inches(7), Inches(1.0), "", font_size=18)
p = gh_tf.paragraphs[0]
p.text = "GitHub Repository"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = CLR_TITLE
add_para(gh_tf, "github.com/MarikSH/Hakaton_SET_ZakupBot", font_size=17, bold=True, color=CLR_ACCENT)
add_para(gh_tf, "Full source code - backend (FastAPI), frontend (React), Docker Compose", font_size=15, color=CLR_SUBTEXT)

# GitHub QR
if os.path.exists(QR_GITHUB_PATH):
    slide5.shapes.add_picture(QR_GITHUB_PATH, Inches(9.5), Inches(1.4), Inches(2.2), Inches(2.2))
add_textbox(slide5, Inches(9.5), Inches(3.65), Inches(2.5), Inches(0.3),
            "Scan to open GitHub", font_size=13, alignment=PP_ALIGN.CENTER, color=CLR_SUBTEXT)

# Deploy link card
deploy_url_text = "[URL after deployment]"
deploy_card = add_rounded_rect(slide5, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.5), CLR_CARD_BG, CLR_GREEN)
deploy_tf = add_textbox(slide5, Inches(1.2), Inches(4.2), Inches(7), Inches(1.0), "", font_size=18)
p = deploy_tf.paragraphs[0]
p.text = "Deployed Product (Latest Version)"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = CLR_TITLE
add_para(deploy_tf, deploy_url_text, font_size=17, bold=True, color=CLR_SUBTEXT)
add_para(deploy_tf, "Deploy: Docker Compose on VPS or Railway (backend) + Vercel (frontend)", font_size=15, color=CLR_SUBTEXT)

# Deploy QR placeholder
if os.path.exists(QR_DEPLOY_PATH):
    slide5.shapes.add_picture(QR_DEPLOY_PATH, Inches(9.5), Inches(4.3), Inches(2.2), Inches(2.2))
add_textbox(slide5, Inches(9.5), Inches(6.55), Inches(2.5), Inches(0.3),
            "Scan after deploy", font_size=13, alignment=PP_ALIGN.CENTER, color=CLR_SUBTEXT)

# Deploy instructions
deploy_note = add_textbox(slide5, Inches(0.8), Inches(6.8), Inches(8.2), Inches(0.5),
            "Deploy steps: 1) Rent VPS / Use Railway free tier  2) docker compose up -d  3) Replace URL above and regenerate QR",
            font_size=13, color=CLR_SUBTEXT)

add_shape(slide5, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), CLR_ACCENT)

add_speaker_notes(slide5,
    "Key links:\n\n"
    "GitHub: github.com/MarikSH/Hakaton_SET_ZakupBot\n"
    "Complete source code for both backend and frontend with Docker Compose.\n\n"
    "Deployed Product: URL placeholder - will be replaced after deployment.\n"
    "Options:\n"
    "- VPS (~500 rub/mo): Install Docker, run docker-compose up\n"
    "- Railway (free): Backend with PostgreSQL\n"
    "- Vercel (free): Frontend static hosting\n\n"
    "QR codes generated via qrserver.com API.\n\n"
    "Thank you! Happy to answer questions.")


# ===== SAVE =====
output_path = os.path.join(SCRIPT_DIR, "Project_Presentation_v2.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
